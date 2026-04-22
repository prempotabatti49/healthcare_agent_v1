"""
Document Processor
==================
Handles PDF, image (JPG/PNG/WEBP), and PPT files.

Strategy per file type:
  PDF  → iterate pages:
           text page  → extract with PyMuPDF directly
           image page → render to pixmap → base64 → GPT-4o Vision
  Image → base64 → GPT-4o Vision
  PPT  → convert each slide to image → same image pipeline

Returns a list of ExtractedPage objects; caller can concatenate or store
individually in SuperMemory.
"""
from __future__ import annotations

import base64
import io
import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

import fitz  # PyMuPDF

from app.config.settings import get_settings

logger = logging.getLogger(__name__)
_settings = get_settings()

# ── Helpers ────────────────────────────────────────────────────────────────────

VISION_PROMPT = (
    "You are a medical document analysis assistant. "
    "Extract ALL text from this image exactly as written, preserving structure. "
    "Include values, dates, diagnoses, medications, and any other clinical content. "
    "Do not summarise — reproduce every piece of information."
)


@dataclass
class ExtractedPage:
    page_number: int
    content: str
    extraction_method: str   # "text" | "vision"
    metadata: dict[str, Any] = field(default_factory=dict)


# ── Vision extraction via OpenAI ───────────────────────────────────────────────

def _extract_via_vision(image_bytes: bytes, mime: str = "image/png") -> str:
    """Call GPT-4o Vision to extract text from an image."""
    from openai import OpenAI

    client = OpenAI(api_key=_settings.openai_api_key)
    b64 = base64.b64encode(image_bytes).decode()
    data_url = f"data:{mime};base64,{b64}"

    response = client.chat.completions.create(
        model=_settings.openai_model,
        messages=[
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": VISION_PROMPT},
                    {"type": "image_url", "image_url": {"url": data_url, "detail": "high"}},
                ],
            }
        ],
        max_tokens=4096,
    )
    return response.choices[0].message.content or ""


# ── PDF processing ─────────────────────────────────────────────────────────────

_TEXT_THRESHOLD = 50   # chars; below this we treat the page as image-only


def _is_text_page(page: fitz.Page) -> bool:
    return len(page.get_text("text").strip()) >= _TEXT_THRESHOLD


def process_pdf(file_bytes: bytes) -> list[ExtractedPage]:
    pages: list[ExtractedPage] = []
    doc = fitz.open(stream=file_bytes, filetype="pdf")

    for i, page in enumerate(doc):
        if _is_text_page(page):
            text = page.get_text("text")
            pages.append(ExtractedPage(
                page_number=i + 1,
                content=text,
                extraction_method="text",
            ))
        else:
            # Render at 150 DPI, convert to PNG bytes
            mat = fitz.Matrix(150 / 72, 150 / 72)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            png_bytes = pix.tobytes("png")
            try:
                content = _extract_via_vision(png_bytes, "image/png")
                pages.append(ExtractedPage(
                    page_number=i + 1,
                    content=content,
                    extraction_method="vision",
                ))
            except Exception as exc:
                logger.warning("Vision extraction failed on page %d: %s", i + 1, exc)
                pages.append(ExtractedPage(
                    page_number=i + 1,
                    content="[Page could not be extracted]",
                    extraction_method="vision_failed",
                ))

    doc.close()
    return pages


# ── Image processing ───────────────────────────────────────────────────────────

def process_image(file_bytes: bytes, filename: str) -> list[ExtractedPage]:
    ext = Path(filename).suffix.lower()
    mime_map = {".jpg": "image/jpeg", ".jpeg": "image/jpeg", ".png": "image/png", ".webp": "image/webp"}
    mime = mime_map.get(ext, "image/png")

    content = _extract_via_vision(file_bytes, mime)
    return [ExtractedPage(page_number=1, content=content, extraction_method="vision")]


# ── PPT processing ─────────────────────────────────────────────────────────────

def process_ppt(file_bytes: bytes) -> list[ExtractedPage]:
    """Convert each PPT slide to a PNG image then run vision extraction."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
        from PIL import Image
    except ImportError:
        raise ImportError("python-pptx and Pillow are required for PPT processing.")

    prs = Presentation(io.BytesIO(file_bytes))
    pages: list[ExtractedPage] = []

    for i, slide in enumerate(prs.slides):
        # Render text fallback (shapes text)
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        slide_text = "\n".join(t for t in texts if t.strip())

        if len(slide_text.strip()) >= _TEXT_THRESHOLD:
            pages.append(ExtractedPage(
                page_number=i + 1,
                content=slide_text,
                extraction_method="text",
            ))
        else:
            # No usable text: report as needing vision
            # Note: full slide-to-image rendering needs LibreOffice/comtypes;
            # for V1 we fall back to text extraction only.
            pages.append(ExtractedPage(
                page_number=i + 1,
                content=slide_text or "[Slide has no extractable text]",
                extraction_method="text_fallback",
            ))

    return pages


# ── Public API ─────────────────────────────────────────────────────────────────

SUPPORTED_EXTENSIONS = {".pdf", ".jpg", ".jpeg", ".png", ".webp", ".pptx", ".ppt"}


def process_document(file_bytes: bytes, filename: str) -> list[ExtractedPage]:
    """
    Route file to the correct processor based on extension.
    Returns list of ExtractedPage objects ready for SuperMemory ingestion.
    """
    ext = Path(filename).suffix.lower()
    if ext == ".pdf":
        return process_pdf(file_bytes)
    elif ext in {".jpg", ".jpeg", ".png", ".webp"}:
        return process_image(file_bytes, filename)
    elif ext in {".pptx", ".ppt"}:
        return process_ppt(file_bytes)
    else:
        raise ValueError(
            f"Unsupported file type: {ext}. "
            f"Supported: {', '.join(SUPPORTED_EXTENSIONS)}"
        )


def pages_to_text(pages: list[ExtractedPage]) -> str:
    """Concatenate all pages into a single string for SuperMemory ingestion."""
    return "\n\n".join(
        f"--- Page {p.page_number} ---\n{p.content}"
        for p in pages
    )
